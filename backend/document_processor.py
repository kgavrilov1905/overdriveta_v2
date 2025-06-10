"""
Document processing module for Alberta Perspectives RAG API
Handles PDF and PowerPoint text extraction, content chunking, and document preprocessing.
"""

import logging
import re
from typing import List, Dict, Any, Tuple, Optional
from datetime import datetime
import hashlib

# Initialize logger first
logger = logging.getLogger(__name__)

# Optional imports with fallbacks
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    logger.warning("PyMuPDF not available. PDF processing will be disabled.")

try:
    from pptx import Presentation  # python-pptx
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False
    logger.warning("python-pptx not available. PowerPoint processing will be disabled.")

class DocumentProcessor:
    """Handles document processing operations."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def extract_text_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text content from PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        if not HAS_PYMUPDF:
            raise ImportError("PyMuPDF is not installed. PDF processing is not available.")
        
        try:
            doc = fitz.open(file_path)
            
            # Extract document metadata
            metadata = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
                "creator": doc.metadata.get("creator", ""),
                "producer": doc.metadata.get("producer", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
                "modification_date": doc.metadata.get("modDate", ""),
                "page_count": doc.page_count
            }
            
            # Extract text from each page
            pages = []
            full_text = ""
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # Extract text
                text = page.get_text()
                
                # Clean and normalize text
                cleaned_text = self._clean_text(text)
                
                if cleaned_text.strip():  # Only add non-empty pages
                    page_info = {
                        "page_number": page_num + 1,
                        "text": cleaned_text,
                        "char_count": len(cleaned_text),
                        "word_count": len(cleaned_text.split())
                    }
                    pages.append(page_info)
                    full_text += f"\n\n--- Page {page_num + 1} ---\n\n{cleaned_text}"
            
            doc.close()
            
            result = {
                "metadata": metadata,
                "pages": pages,
                "full_text": full_text.strip(),
                "total_pages": len(pages),
                "total_chars": len(full_text),
                "total_words": len(full_text.split())
            }
            
            logger.info(f"Extracted text from PDF: {len(pages)} pages, {len(full_text)} characters")
            return result
            
        except Exception as e:
            logger.error(f"Failed to extract text from PDF {file_path}: {str(e)}")
            raise
    
    def extract_text_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text content from PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        if not HAS_PYTHON_PPTX:
            raise ImportError("python-pptx is not installed. PowerPoint processing is not available.")
        
        try:
            prs = Presentation(file_path)
            
            # Extract document metadata
            core_props = prs.core_properties
            metadata = {
                "title": getattr(core_props, 'title', '') or "",
                "author": getattr(core_props, 'author', '') or "",
                "subject": getattr(core_props, 'subject', '') or "",
                "creator": getattr(core_props, 'author', '') or "",  # PowerPoint uses author as creator
                "creation_date": getattr(core_props, 'created', '') or "",
                "modification_date": getattr(core_props, 'modified', '') or "",
                "slide_count": len(prs.slides)
            }
            
            # Extract text from each slide
            slides = []
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Handle table text
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text += " | ".join(row_text) + "\n"
                
                # Clean and normalize text
                cleaned_text = self._clean_text(slide_text)
                
                if cleaned_text.strip():  # Only add non-empty slides
                    slide_info = {
                        "page_number": slide_num + 1,  # Using page_number for consistency with PDF
                        "text": cleaned_text,
                        "char_count": len(cleaned_text),
                        "word_count": len(cleaned_text.split())
                    }
                    slides.append(slide_info)
                    full_text += f"\n\n--- Slide {slide_num + 1} ---\n\n{cleaned_text}"
            
            result = {
                "metadata": metadata,
                "pages": slides,  # Using 'pages' for consistency with PDF processing
                "full_text": full_text.strip(),
                "total_pages": len(slides),
                "total_chars": len(full_text),
                "total_words": len(full_text.split())
            }
            
            logger.info(f"Extracted text from PowerPoint: {len(slides)} slides, {len(full_text)} characters")
            return result
            
        except Exception as e:
            logger.error(f"Failed to extract text from PowerPoint {file_path}: {str(e)}")
            raise
    
    def _clean_text(self, text: str) -> str:
        """
        Clean and normalize extracted text.
        
        Args:
            text: Raw text from PDF
            
        Returns:
            Cleaned text
        """
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove page headers/footers patterns (common in reports)
        text = re.sub(r'^\d+\s*$', '', text, flags=re.MULTILINE)  # Page numbers on their own line
        text = re.sub(r'^Page \d+ of \d+.*$', '', text, flags=re.MULTILINE)  # Page X of Y
        
        # Clean up common PDF artifacts
        text = re.sub(r'(?i)continued on next page', '', text)
        text = re.sub(r'(?i)see page \d+', '', text)
        
        # Normalize quotes and dashes
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace(''', "'").replace(''', "'")
        text = text.replace('–', '-').replace('—', '-')
        
        # Remove excessive line breaks
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        
        return text.strip()
    
    def chunk_text(self, text: str, page_number: Optional[int] = None) -> List[Dict[str, Any]]:
        """
        Split text into chunks for embedding generation.
        
        Args:
            text: Text to chunk
            page_number: Page number for reference
            
        Returns:
            List of text chunks with metadata
        """
        if not text.strip():
            return []
        
        chunks = []
        
        # Split by sentences first
        sentences = self._split_into_sentences(text)
        
        current_chunk = ""
        current_chunk_sentences = []
        chunk_index = 0
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) > self.chunk_size and current_chunk:
                # Create chunk
                chunk_info = self._create_chunk_info(
                    content=current_chunk.strip(),
                    chunk_index=chunk_index,
                    page_number=page_number,
                    sentences=current_chunk_sentences
                )
                chunks.append(chunk_info)
                
                # Start new chunk with overlap
                overlap_sentences = current_chunk_sentences[-self._calculate_overlap_sentences(current_chunk_sentences):]
                current_chunk = " ".join(overlap_sentences)
                current_chunk_sentences = overlap_sentences[:]
                chunk_index += 1
            
            current_chunk += " " + sentence if current_chunk else sentence
            current_chunk_sentences.append(sentence)
        
        # Add final chunk if it has content
        if current_chunk.strip():
            chunk_info = self._create_chunk_info(
                content=current_chunk.strip(),
                chunk_index=chunk_index,
                page_number=page_number,
                sentences=current_chunk_sentences
            )
            chunks.append(chunk_info)
        
        logger.info(f"Created {len(chunks)} chunks from text (page {page_number or 'N/A'})")
        return chunks
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences."""
        # Simple sentence splitting - could be improved with spaCy or NLTK
        sentences = re.split(r'(?<=[.!?])\s+', text)
        return [s.strip() for s in sentences if s.strip()]
    
    def _calculate_overlap_sentences(self, sentences: List[str]) -> int:
        """Calculate number of sentences to include in overlap."""
        total_overlap_chars = 0
        overlap_sentences = 0
        
        for sentence in reversed(sentences):
            if total_overlap_chars + len(sentence) <= self.chunk_overlap:
                total_overlap_chars += len(sentence)
                overlap_sentences += 1
            else:
                break
        
        return min(overlap_sentences, len(sentences) // 2)  # Max 50% overlap
    
    def _create_chunk_info(self, content: str, chunk_index: int, page_number: Optional[int], sentences: List[str]) -> Dict[str, Any]:
        """Create chunk information dictionary."""
        return {
            "content": content,
            "chunk_index": chunk_index,
            "page_number": page_number,
            "char_count": len(content),
            "word_count": len(content.split()),
            "sentence_count": len(sentences),
            "content_hash": hashlib.md5(content.encode()).hexdigest()
        }
    
    def process_document(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """
        Complete document processing pipeline.
        
        Args:
            file_path: Path to the document file
            file_name: Original filename
            
        Returns:
            Processed document data ready for database storage
        """
        logger.info(f"Starting document processing for: {file_name}")
        
        try:
            # Determine file type and extract text
            file_extension = file_name.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                extracted_data = self.extract_text_from_pdf(file_path)
                content_type = "application/pdf"
            elif file_extension in ['pptx', 'ppt']:
                extracted_data = self.extract_text_from_pptx(file_path)
                content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                raise ValueError(f"Unsupported file type: {file_extension}. Supported types: PDF, PowerPoint (PPTX)")
            
            # Process chunks for each page/slide
            all_chunks = []
            chunk_index = 0
            
            for page_info in extracted_data["pages"]:
                page_chunks = self.chunk_text(
                    text=page_info["text"],
                    page_number=page_info["page_number"]
                )
                
                # Update global chunk indices
                for chunk in page_chunks:
                    chunk["chunk_index"] = chunk_index
                    chunk_index += 1
                
                all_chunks.extend(page_chunks)
            
            # Also create chunks from full text for cross-page content
            full_text_chunks = self.chunk_text(extracted_data["full_text"])
            for chunk in full_text_chunks:
                chunk["chunk_index"] = chunk_index
                chunk["page_number"] = None  # Cross-page chunk
                chunk_index += 1
            
            all_chunks.extend(full_text_chunks)
            
            # Create document metadata
            document_data = {
                "file_name": file_name,
                "title": extracted_data["metadata"]["title"] or file_name,
                "content_type": content_type,
                "page_count": extracted_data["total_pages"],
                "processing_status": "processing",
                "metadata": {
                    **extracted_data["metadata"],
                    "total_chunks": len(all_chunks),
                    "processing_timestamp": datetime.utcnow().isoformat(),
                    "chunk_size": self.chunk_size,
                    "chunk_overlap": self.chunk_overlap,
                    "file_type": file_extension.upper()
                }
            }
            
            result = {
                "document": document_data,
                "chunks": all_chunks,
                "extracted_text": extracted_data
            }
            
            logger.info(f"Document processing completed: {len(all_chunks)} chunks created")
            return result
            
        except Exception as e:
            logger.error(f"Document processing failed for {file_name}: {str(e)}")
            raise 